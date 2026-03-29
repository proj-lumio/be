"""Document ingestion pipeline — MongoDB version."""

import asyncio
import io
import logging
import os
import tempfile
from datetime import datetime, timezone
from pathlib import Path

from bson import ObjectId

from app.config import get_settings
from app.db.mongo import get_db

logger = logging.getLogger(__name__)
settings = get_settings()

INGESTION_CONCURRENCY = 10  # max parallel document processing tasks


# ── Text extractors ──

def extract_pdf(b: bytes) -> str:
    from PyPDF2 import PdfReader
    return "\n".join(p.extract_text() or "" for p in PdfReader(io.BytesIO(b)).pages)

def extract_docx(b: bytes) -> str:
    from docx import Document
    return "\n".join(p.text for p in Document(io.BytesIO(b)).paragraphs)

def extract_xlsx(b: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(b), read_only=True)
    lines = []
    for s in wb.sheetnames:
        for row in wb[s].iter_rows(values_only=True):
            lines.append(" | ".join(str(c) if c else "" for c in row))
    return "\n".join(lines)

def extract_pptx(b: bytes) -> str:
    from pptx import Presentation
    texts = []
    for slide in Presentation(io.BytesIO(b)).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text)
    return "\n".join(texts)

EXTRACTORS = {
    "pdf": extract_pdf, "docx": extract_docx,
    "xlsx": extract_xlsx, "pptx": extract_pptx,
    "txt": lambda b: b.decode("utf-8", errors="replace"),
}


def detect_doc_type(filename: str) -> str:
    ext = Path(filename).suffix.lower()
    mapping = {
        ".pdf": "pdf", ".docx": "docx", ".doc": "docx",
        ".xlsx": "xlsx", ".xls": "xlsx", ".pptx": "pptx", ".ppt": "pptx",
        ".txt": "txt", ".csv": "txt",
        ".mp3": "audio", ".wav": "audio", ".m4a": "audio", ".ogg": "audio",
        ".mp4": "video", ".webm": "video", ".mov": "video", ".avi": "video",
        ".png": "image", ".jpg": "image", ".jpeg": "image",
    }
    return mapping.get(ext, "txt")


def chunk_text(text: str, size: int = 1000, overlap: int = 200) -> list[dict]:
    chunks, start, idx = [], 0, 0
    while start < len(text):
        content = text[start:start + size]
        if content.strip():
            chunks.append({"chunk_index": idx, "content": content})
            idx += 1
        start += size - overlap
    return chunks


async def process_document(doc_id: str, file_bytes: bytes):
    """Full pipeline: extract → chunk → embed (Qdrant) → graph (Neo4j)."""
    db = get_db()
    oid = ObjectId(doc_id)

    doc = await db.documents.find_one({"_id": oid})
    if not doc:
        return

    await db.documents.update_one({"_id": oid}, {"$set": {"processing_status": "processing"}})

    try:
        # 1. Extract text
        doc_type = doc["doc_type"]
        if doc_type in ("audio", "video"):
            suffix = Path(doc["filename"]).suffix
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            try:
                from app.services.speech_to_text import transcribe_audio
                raw_text = await transcribe_audio(tmp_path)
            finally:
                os.unlink(tmp_path)
        elif doc_type == "image":
            from app.services.ocr import extract_text_from_image
            raw_text = await extract_text_from_image(file_bytes)
        else:
            extractor = EXTRACTORS.get(doc_type)
            if not extractor:
                raise ValueError(f"Unsupported type: {doc_type}")
            raw_text = extractor(file_bytes)

        logger.info(f"Extracted {len(raw_text)} chars from {doc['filename']}")

        # 2. Chunk
        chunks = chunk_text(raw_text)
        company_id = str(doc["company_id"])

        # 3. Embed in Qdrant
        point_ids = [None] * len(chunks)
        if settings.qdrant_url and chunks:
            try:
                from app.services.vector_store import upsert_chunks
                point_ids = await upsert_chunks(chunks, company_id, doc_id)
                logger.info(f"Stored {len(point_ids)} vectors in Qdrant")
            except Exception as e:
                logger.warning(f"Qdrant failed: {e}")

        # 4. Store chunks in MongoDB
        if chunks:
            chunk_docs = [
                {
                    "document_id": oid,
                    "company_id": doc["company_id"],
                    "chunk_index": c["chunk_index"],
                    "content": c["content"],
                    "qdrant_point_id": pid,
                    "created_at": datetime.now(timezone.utc),
                }
                for c, pid in zip(chunks, point_ids)
            ]
            await db.document_chunks.insert_many(chunk_docs)

        # 5. GraphRAG entities (Neo4j)
        if settings.neo4j_uri and raw_text:
            try:
                from app.services.llm import extract_entities
                from app.services.graph_store import store_entities_and_relations
                extracted = await extract_entities(raw_text[:8000])
                entities = extracted.get("entities", [])
                relationships = extracted.get("relationships", [])
                if entities:
                    store_entities_and_relations(company_id, doc_id, entities, relationships)
            except Exception as e:
                logger.warning(f"GraphRAG failed: {e}")

        # 5b. Category classification
        if settings.neo4j_uri and raw_text:
            try:
                from app.services.llm import classify_company_categories
                from app.services.graph_store import store_company_categories

                categories = await classify_company_categories(raw_text[:6000])
                if categories:
                    store_company_categories(company_id, categories)
                    await db.companies.update_one(
                        {"_id": doc["company_id"]},
                        {"$addToSet": {"categories": {"$each": categories}}},
                    )
                    logger.info(f"Assigned categories {categories} to company {company_id}")
            except Exception as e:
                logger.warning(f"Category classification failed: {e}")

        # 6. Contract analysis (structured extraction)
        if raw_text:
            try:
                from app.services.llm import extract_contract_data
                contract = await extract_contract_data(raw_text)
                if contract:
                    contract.update({
                        "document_id": oid,
                        "company_id": doc["company_id"],
                        "created_at": datetime.now(timezone.utc),
                        "updated_at": datetime.now(timezone.utc),
                        "criticality_manual": None,
                    })
                    await db.contract_analyses.update_one(
                        {"document_id": oid},
                        {"$set": contract},
                        upsert=True,
                    )
                    logger.info(f"Contract analysis saved for {doc['filename']}")
            except Exception as e:
                logger.warning(f"Contract analysis failed: {e}")

        await db.documents.update_one({"_id": oid}, {"$set": {
            "raw_text": raw_text,
            "processing_status": "completed",
            "updated_at": datetime.now(timezone.utc),
        }})
        logger.info(f"Document {doc['filename']} processed OK")

    except Exception as e:
        await db.documents.update_one({"_id": oid}, {"$set": {
            "processing_status": "failed",
            "error_message": str(e)[:500],
            "updated_at": datetime.now(timezone.utc),
        }})
        logger.exception(f"Processing failed: {e}")
        raise


async def process_documents_batch(doc_ids_and_bytes: list[tuple[str, bytes]]):
    """Process multiple documents in parallel with concurrency control."""
    sem = asyncio.Semaphore(INGESTION_CONCURRENCY)

    async def _guarded(doc_id: str, file_bytes: bytes):
        async with sem:
            await process_document(doc_id, file_bytes)

    await asyncio.gather(
        *[_guarded(did, fb) for did, fb in doc_ids_and_bytes],
        return_exceptions=True,
    )
